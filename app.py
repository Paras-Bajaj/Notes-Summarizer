from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import os
import uuid
from werkzeug.utils import secure_filename
import PyPDF2
from pptx import Presentation
import pytesseract
from PIL import Image
import time
import re
import nltk

# Download NLTK punkt if not already present
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Initialize Flask app
app = Flask(__name__)
CORS(app)  # Enable CORS for all origins

# Configuration
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
app.config['SECRET_KEY'] = os.urandom(24).hex()

# Ensure upload directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Allowed file extensions
ALLOWED_EXTENSIONS = {
    'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 
    'ppt', 'pptx', 'doc', 'docx'
}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    return text

def extract_text_from_ppt(file_path):
    """Extract text from PowerPoint file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPT: {e}")
    return text

def extract_text_from_image(file_path):
    """Extract text from image using OCR"""
    try:
        # Try using pytesseract first
        try:
            image = Image.open(file_path)
            if image.mode != "RGB":
                image = image.convert("RGB")
            text = pytesseract.image_to_string(image, lang='eng')
            print(f"OCR extracted text (tesseract): {text[:100]}")
            return text
        except Exception as e:
            print(f"Tesseract failed: {e}")
            
        # Fallback to easyocr if available
        try:
            import easyocr
            reader = easyocr.Reader(['en'])
            result = reader.readtext(file_path, detail=0, paragraph=True)
            text = ' '.join(result)
            print(f"OCR extracted text (easyocr): {text[:100]}")
            return text
        except ImportError:
            print("EasyOCR not installed")
            
        return "OCR error: Please install Tesseract OCR or EasyOCR for image text extraction."
        
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return "Error processing image. Please try another file."

def extract_text_from_txt(file_path):
    """Extract text from TXT file"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"Error extracting text from TXT: {e}")
        return ""

def extract_text_from_docx(file_path):
    """Extract text from DOCX file"""
    try:
        import docx2txt
        text = docx2txt.process(file_path)
        return text
    except ImportError:
        return "Please install docx2txt to process Word documents: pip install docx2txt"
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""

def simple_summarize(text, max_sentences=3):
    """Simple extractive summarization by sentence ranking"""
    if not text or not text.strip():
        return "No content to summarize."
    
    # Tokenize sentences using NLTK
    try:
        sentences = nltk.sent_tokenize(text)
        if len(sentences) <= 1:
            return text[:200] + "..." if len(text) > 200 else text
            
        # Frequency-based summarization
        words = [w for w in nltk.word_tokenize(text.lower()) if w.isalnum()]
        if not words:
            return text[:200] + "..." if len(text) > 200 else text
            
        freq = nltk.FreqDist(words)
        
        # Score sentences
        ranked_sentences = sorted(
            sentences,
            key=lambda s: sum(freq[w.lower()] for w in nltk.word_tokenize(s) if w.isalnum()),
            reverse=True
        )
        
        # Get top sentences while preserving order
        important_sentences = []
        for sentence in sentences:
            if sentence in ranked_sentences[:max_sentences]:
                important_sentences.append(sentence)
                
        summary = ' '.join(important_sentences)
        return summary if summary.strip() else text[:200] + "..." if len(text) > 200 else text
        
    except Exception as e:
        print(f"Error in summarization: {e}")
        # Fallback: return first few sentences
        sentences = text.split('.')
        return '.'.join(sentences[:max_sentences]) + '.'

@app.route('/')
def serve_frontend():
    return send_from_directory('.', 'intro.html')

# Serve static files (CSS, JS, etc.)
@app.route('/<path:path>')
def serve_static(path):
    return send_from_directory('.', path)

@app.route('/api/summarize', methods=['POST'])
def summarize():
    start_time = time.time()
    print("Received summarize request")

    # Handle CORS preflight
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'ok'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
        response.headers.add('Access-Control-Allow-Methods', 'POST, OPTIONS')
        return response
    
    text = None
    if 'text' in request.form:
        text = request.form['text']
    elif request.is_json:
        data = request.get_json()
        text = data.get('text')

    if text is not None:
        print(f"Text input received: {str(text)[:100]}")
        summary = simple_summarize(text)
        processing_time = int((time.time() - start_time) * 1000)
        original_length = len(text.split())
        summary_length = len(summary.split())
        compression_ratio = int((1 - (summary_length / original_length)) * 100) if original_length > 0 else 0

        response = jsonify({
            'summary': summary if summary.strip() else "No summary generated.",
            'processing_time': processing_time,
            'compression_ratio': compression_ratio,
            'original_length': original_length,
            'summary_length': summary_length
        })
        response.headers.add('Access-Control-Allow-Origin', '*')
        return response

    elif 'file' in request.files:
        file = request.files['file']
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            unique_filename = f"{uuid.uuid4().hex}_{filename}"
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            file.save(file_path)
            file_ext = filename.rsplit('.', 1)[1].lower()
            text = ""

            print(f"File upload received: {filename}")

            try:
                if file_ext == 'pdf':
                    text = extract_text_from_pdf(file_path)
                elif file_ext in ['ppt', 'pptx']:
                    text = extract_text_from_ppt(file_path)
                elif file_ext in ['png', 'jpg', 'jpeg', 'gif']:
                    text = extract_text_from_image(file_path)
                elif file_ext == 'txt':
                    text = extract_text_from_txt(file_path)
                elif file_ext in ['doc', 'docx']:
                    text = extract_text_from_docx(file_path)
                else:
                    text = f"Unsupported file type: {file_ext}"
            except Exception as e:
                print(f"Error during file extraction: {e}")
                if os.path.exists(file_path):
                    os.remove(file_path)
                response = jsonify({'error': f'Error extracting text: {str(e)}'})
                response.headers.add('Access-Control-Allow-Origin', '*')
                return response, 500

            if os.path.exists(file_path):
                os.remove(file_path)
                
            print(f"Extracted text: {text[:100]}")

            # If no text, return a user-friendly message
            if not text or not text.strip():
                response = jsonify({
                    'summary': "No text detected. Please try a different file or check if OCR is properly installed.",
                    'processing_time': int((time.time() - start_time) * 1000),
                    'compression_ratio': 0,
                    'original_length': 0,
                    'summary_length': 0
                })
                response.headers.add('Access-Control-Allow-Origin', '*')
                return response

            summary = simple_summarize(text)
            processing_time = int((time.time() - start_time) * 1000)
            original_length = len(text.split())
            summary_length = len(summary.split())
            compression_ratio = int((1 - (summary_length / original_length)) * 100) if original_length > 0 else 0

            response = jsonify({
                'summary': summary if summary.strip() else "No summary generated.",
                'processing_time': processing_time,
                'compression_ratio': compression_ratio,
                'original_length': original_length,
                'summary_length': summary_length
            })
            response.headers.add('Access-Control-Allow-Origin', '*')
            return response
        else:
            response = jsonify({'error': 'Invalid file type. Supported formats: PDF, PPT, DOC, TXT, and images.'})
            response.headers.add('Access-Control-Allow-Origin', '*')
            return response, 400

    else:
        print("No text or file provided in request")
        response = jsonify({'error': 'No text or file provided'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        return response, 400

@app.route('/api/sample', methods=['GET'])
def get_sample_text():
    """Return sample text for demonstration"""
    samples = [
        "Artificial intelligence (AI) is revolutionizing the way we work, live, and interact with technology. From machine learning algorithms that can predict consumer behavior to natural language processing systems that can understand and respond to human speech, AI is transforming industries across the globe. In healthcare, AI is being used to diagnose diseases more accurately and develop personalized treatment plans. In finance, AI algorithms are detecting fraud and making investment decisions. In transportation, autonomous vehicles powered by AI are becoming a reality. As AI continues to evolve, it promises to bring even more innovative solutions to complex problems, making our lives more efficient and productive.",
        "Climate change represents one of the most pressing challenges of our time, with far-reaching implications for ecosystems, human societies, and the global economy. Rising global temperatures, caused primarily by greenhouse gas emissions from human activities, are leading to more frequent and severe weather events, including hurricanes, droughts, and floods. The melting of polar ice caps and glaciers is contributing to rising sea levels, threatening coastal communities worldwide. To address this crisis, governments, businesses, and individuals must work together to reduce carbon emissions, transition to renewable energy sources, and implement sustainable practices. The Paris Agreement represents a significant step forward in global climate action, but much more needs to be done to limit global warming and protect our planet for future generations."
    ]
    response = jsonify({'sample': samples[0]})
    response.headers.add('Access-Control-Allow-Origin', '*')
    return response

if __name__ == '__main__':
    app.run(debug=True, port=5000, host='0.0.0.0')